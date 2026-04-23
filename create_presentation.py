from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Design tokens ──────────────────────────────────────────────────────────────
BLUE       = RGBColor(0x1E, 0x50, 0xA2)   # main brand blue
BLUE_LIGHT = RGBColor(0xD6, 0xE4, 0xF7)   # pale blue fill
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1A, 0x1A, 0x2E)
GRAY       = RGBColor(0x55, 0x55, 0x55)
ACCENT     = RGBColor(0xE8, 0x4A, 0x1E)   # orange-red accent

W = Inches(13.33)   # slide width  (16:9 widescreen)
H = Inches(7.50)    # slide height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helper functions ───────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def add_bullets(slide, items, x, y, w, h, size=17, title_size=None,
                color=DARK, indent_color=BLUE, spacing_after=120000):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Emu(spacing_after)
        if isinstance(item, tuple):
            label, body = item
            run1 = p.add_run()
            run1.text = f"•  {label}: "
            run1.font.size  = Pt(size)
            run1.font.bold  = True
            run1.font.color.rgb = indent_color
            run2 = p.add_run()
            run2.text = body
            run2.font.size  = Pt(size)
            run2.font.bold  = False
            run2.font.color.rgb = color
        else:
            run = p.add_run()
            run.text = f"•  {item}"
            run.font.size  = Pt(size)
            run.font.bold  = False
            run.font.color.rgb = color
    return txb

def slide_header(slide, title, subtitle=None, bar_h=1.05):
    """Blue top bar with white title text."""
    add_rect(slide, 0, 0, 13.33, bar_h, BLUE)
    add_text(slide, title,
             x=0.35, y=0.10, w=12.0, h=bar_h - 0.1,
             size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 x=0.35, y=bar_h - 0.08, w=12.0, h=0.38,
                 size=13, bold=False, color=BLUE_LIGHT, italic=True)
    # thin accent line below bar
    add_rect(slide, 0, bar_h, 13.33, 0.04, ACCENT)

def add_image(slide, path, x, y, w, h=None):
    if not os.path.exists(path):
        return
    if h:
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))

def stat_box(slide, label, value, x, y, w=2.8, h=1.15):
    add_rect(slide, x, y, w, h, BLUE_LIGHT)
    add_rect(slide, x, y, w, 0.06, BLUE)
    add_text(slide, value, x+0.1, y+0.08, w-0.2, 0.55,
             size=26, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x+0.1, y+0.62, w-0.2, 0.45,
             size=11, bold=False, color=GRAY, align=PP_ALIGN.CENTER)

def footer(slide, text="Sama Gannurdinov  |  Data Mining Final Project  |  April 2026"):
    add_rect(slide, 0, 7.30, 13.33, 0.20, RGBColor(0xF0, 0xF4, 0xFB))
    add_text(slide, text, 0.3, 7.31, 12.7, 0.18,
             size=9, color=GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.50, BLUE)                        # full blue bg
add_rect(sl, 0, 4.80, 13.33, 2.70, RGBColor(0x14, 0x36, 0x78))  # darker lower band
add_rect(sl, 0, 4.77, 13.33, 0.06, ACCENT)                   # accent divider

add_text(sl, "Chicago Crime Dataset 2024",
         0.6, 1.2, 12.0, 1.0, size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Data Mining Final Project",
         0.6, 2.35, 12.0, 0.8, size=28, bold=False, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

add_rect(sl, 4.5, 3.35, 4.33, 0.05, ACCENT)

add_text(sl, "Preprocessing  ·  Data Integration  ·  Cluster Analysis  ·  Classification",
         0.6, 3.55, 12.0, 0.55, size=15, bold=False, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

add_text(sl, "Sama Gannurdinov", 0.6, 5.05, 12.0, 0.45,
         size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Data Mining  ·  April 2026", 0.6, 5.55, 12.0, 0.40,
         size=13, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
add_text(sl, "github.com/samagannurdinov", 0.6, 6.10, 12.0, 0.35,
         size=12, color=BLUE_LIGHT, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT'S NEW (Midterm → Final)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Project Overview", "What was added since the Midterm")
footer(sl)

# Two column panels
for col_x, col_title, col_color, items in [
    (0.35, "Midterm  (completed)", RGBColor(0x4A, 0x4A, 0x8A), [
        "Data Cleaning  — 99.4% records retained",
        "Attribute Subset Selection  — 7 columns removed",
        "PCA  — 84.21% variance in 2 components",
        "Min-Max & Z-Score Normalization",
        "Equal-Width Binning (time periods)",
        "Concept Hierarchy (77 areas → 4 sectors)",
        "Binary Discretization (indoor / outdoor)",
        "Basic K-Means Clustering  (Lat / Long only)",
        "Interactive D3.js Dashboard",
    ]),
    (7.00, "Final  (new progress)", BLUE, [
        "Data Integration  — 2020 Census demographics merged for all 77 community areas",
        "Per-capita crime rates computed by neighborhood",
        "Enhanced Clustering  — Elbow Method validates K=10",
        "4-feature K-Means (location + time + environment)",
        "Cluster profiling  — arrest rate & avg hour per cluster",
        "Classification Model  — Random Forest (arrest prediction)",
        "84.3% accuracy  ·  ROC-AUC 0.8154",
        "Feature Importance analysis",
        "8-section Final PDF Report with 10 charts",
    ]),
]:
    add_rect(sl, col_x, 1.20, 5.85, 5.85, RGBColor(0xF5, 0xF8, 0xFF))
    add_rect(sl, col_x, 1.20, 5.85, 0.42, col_color)
    add_text(sl, col_title, col_x + 0.15, 1.22, 5.5, 0.38,
             size=14, bold=True, color=WHITE)
    add_bullets(sl, items, col_x + 0.20, 1.72, 5.50, 5.10,
                size=13, color=DARK, indent_color=col_color, spacing_after=60000)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DATASET STATS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Dataset Overview", "City of Chicago Open Data Portal — 2024")
footer(sl)

# 6 stat boxes in 2 rows
stats = [
    ("Raw Records",       "259,032"),
    ("After Cleaning",    "257,547"),
    ("Retention Rate",    "99.4%"),
    ("Original Features", "22"),
    ("Engineered Cols",   "28"),
    ("Community Areas",   "77"),
]
for i, (label, val) in enumerate(stats):
    col = i % 3
    row = i // 3
    stat_box(sl, label, val, x=0.40 + col * 3.05, y=1.30 + row * 1.45)

add_text(sl, "Top 5 Crime Types (2024)", 0.40, 4.20, 6.0, 0.40,
         size=15, bold=True, color=BLUE)
crimes = [
    ("Theft",               "60,298",  "23.4%"),
    ("Battery",             "45,953",  "17.8%"),
    ("Criminal Damage",     "28,449",  "11.0%"),
    ("Assault",             "23,380",   "9.1%"),
    ("Motor Vehicle Theft", "21,642",   "8.4%"),
]
# Mini table
hdrs = ["Crime Type", "Count", "%"]
col_ws = [3.4, 1.2, 1.0]
col_xs = [0.40, 3.80, 5.00]
# header row
for hdr, cx, cw in zip(hdrs, col_xs, col_ws):
    add_rect(sl, cx, 4.65, cw - 0.05, 0.36, BLUE)
    add_text(sl, hdr, cx + 0.07, 4.67, cw - 0.12, 0.32,
             size=12, bold=True, color=WHITE)
for ri, (ctype, cnt, pct) in enumerate(crimes):
    bg = BLUE_LIGHT if ri % 2 == 0 else WHITE
    for val, cx, cw in zip([ctype, cnt, pct], col_xs, col_ws):
        add_rect(sl, cx, 5.05 + ri * 0.36, cw - 0.05, 0.34, bg)
        add_text(sl, val, cx + 0.07, 5.07 + ri * 0.36, cw - 0.12, 0.30,
                 size=11, color=DARK)

add_text(sl, "Theft alone accounts for nearly 1 in 4 reported incidents across the city.",
         6.30, 4.20, 6.70, 0.45, size=13, bold=False, color=GRAY, italic=True)
add_image(sl, "correlation_heatmap.png", 6.30, 4.65, 6.60)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PREPROCESSING PIPELINE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Data Preprocessing Pipeline", "9 techniques applied before any modeling")
footer(sl)

steps = [
    ("Data Cleaning",         "Dropped 1,486 rows with missing coordinates; filled Location Description gaps"),
    ("Attribute Selection",   "Removed 7 irrelevant ID columns (IUCR, FBI Code, Case Number…)"),
    ("PCA",                   "4 spatial features → 2 principal components retaining 84.21% variance"),
    ("Min-Max Normalization", "Lat / Long scaled to [0,1] — required for K-Means distance calculation"),
    ("Z-Score Normalization", "District / Ward standardized — required for PCA (mean=0, std=1)"),
    ("Feature Engineering",   "Extracted Hour, Month, Day_of_Week from raw timestamp"),
    ("Equal-Width Binning",   "24 hour values → 4 time periods: Late Night / Morning / Afternoon / Evening"),
    ("Concept Hierarchy",     "77 Community Areas → 4 City Sectors (North Side, Central/West, South, Far South)"),
    ("Binary Discretization", "Dozens of location types → Is_Outdoor flag (1 = outdoor, 0 = indoor)"),
]
for i, (tech, desc) in enumerate(steps):
    row = i // 3
    col = i % 3
    bx = 0.25 + col * 4.36
    by = 1.18 + row * 1.90
    add_rect(sl, bx, by, 4.20, 1.75, BLUE_LIGHT)
    add_rect(sl, bx, by, 4.20, 0.08, BLUE if i % 2 == 0 else ACCENT)
    add_text(sl, tech, bx + 0.12, by + 0.12, 4.0, 0.40,
             size=13, bold=True, color=BLUE)
    add_text(sl, desc, bx + 0.12, by + 0.52, 3.95, 1.15,
             size=11, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PCA RESULT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Data Reduction — PCA", "Compressing 4 correlated spatial features into 2 components")
footer(sl)

add_image(sl, "pca_reduction_plot.png", 0.30, 1.15, 6.30)

add_text(sl, "Why PCA?", 7.0, 1.20, 5.90, 0.45, size=16, bold=True, color=BLUE)
add_bullets(sl, [
    ("District vs Ward correlation", "r = 0.653  — they carry overlapping information"),
    ("PC1 captures", "69.75% of variance  (north-south city axis)"),
    ("PC2 captures", "14.46% of variance  (east-west spread)"),
    ("Total retained", "84.21% of spatial information in just 2 dimensions"),
    ("Pre-processing", "Z-Score normalization applied before PCA to equalize feature scales"),
], 7.0, 1.75, 6.10, 4.80, size=13, indent_color=BLUE)

add_rect(sl, 7.0, 6.30, 6.10, 0.60, RGBColor(0xE8, 0xF0, 0xFF))
add_text(sl, "Result: 50% dimensionality reduction (4 → 2) with minimal information loss.",
         7.10, 6.33, 5.90, 0.54, size=13, bold=True, color=BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DATA INTEGRATION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Data Integration", "Merging 2020 U.S. Census demographics for all 77 Chicago community areas")
footer(sl)

add_image(sl, "crime_rate_by_area.png", 0.25, 1.15, 7.20)

add_text(sl, "Why Integrate?", 7.70, 1.20, 5.30, 0.45, size=16, bold=True, color=BLUE)
add_bullets(sl, [
    "Raw crime counts favor large-population areas — unfair comparison",
    "Per-capita rate reveals the true risk exposure for residents",
    "Enables evidence-based, equity-aware policy decisions",
], 7.70, 1.75, 5.30, 1.40, size=13, spacing_after=70000)

add_text(sl, "Key Findings", 7.70, 3.30, 5.30, 0.45, size=16, bold=True, color=BLUE)
findings = [
    ("Fuller Park",    "243 / 1k residents  — highest in the city"),
    ("Edison Park",    "24 / 1k residents  — lowest in the city"),
    ("10×  gap",       "between safest and most dangerous areas"),
    ("The Loop",       "High raw count (9,289) but explained by large transient population"),
]
add_bullets(sl, findings, 7.70, 3.85, 5.30, 3.0, size=13, indent_color=ACCENT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — POPULATION vs CRIME SCATTER
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Data Integration — Population vs Crime", "Color shows per-capita crime rate — not all busy areas are high-risk")
footer(sl)

add_image(sl, "population_vs_crime.png", 0.30, 1.15, 8.50)

add_text(sl, "Insight", 9.10, 1.20, 3.95, 0.45, size=16, bold=True, color=BLUE)
add_bullets(sl, [
    "Large population does not equal high crime rate",
    "West Garfield Park & Englewood: small population, extreme per-capita rates",
    "Austin: largest population (100k) but mid-tier crime rate",
    "Integration reveals inequality invisible in raw counts",
], 9.10, 1.75, 3.95, 4.50, size=13, spacing_after=90000)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ELBOW METHOD
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Model 1: Cluster Analysis — Elbow Method", "Empirically validating K=10 instead of choosing it arbitrarily")
footer(sl)

add_image(sl, "elbow_curve.png", 0.30, 1.15, 8.10)

add_text(sl, "What is the Elbow Method?", 8.65, 1.20, 4.40, 0.45, size=15, bold=True, color=BLUE)
add_bullets(sl, [
    "Plot inertia (WCSS) for K = 2 to 15",
    "Find the 'elbow' where adding more clusters gives diminishing returns",
    "Inflection occurs near K = 9–10",
    "K = 10 selected as optimal balance between granularity and efficiency",
], 8.65, 1.75, 4.40, 2.70, size=13, spacing_after=80000)

add_text(sl, "Improvement vs Midterm", 8.65, 4.65, 4.40, 0.45, size=15, bold=True, color=BLUE)
add_bullets(sl, [
    ("Midterm", "K=10 chosen with no justification"),
    ("Final",   "K=10 validated by Elbow Method inertia analysis"),
    ("Features", "Added Hour & Is_Outdoor so clusters reflect WHEN and WHERE crimes happen — not just coordinates"),
], 8.65, 5.18, 4.40, 2.10, size=12, indent_color=ACCENT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CLUSTER MAP + PROFILES
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Model 1: Crime Hotspot Clusters (K=10)", "Enhanced K-Means — Location + Time-of-Day + Indoor/Outdoor")
footer(sl)

add_image(sl, "crime_clusters.png", 0.25, 1.12, 5.10)
add_image(sl, "cluster_profiles.png", 5.50, 1.12, 7.55)

add_text(sl, "Outdoor evening clusters (3, 8) have the highest arrest rates: 19–21%",
         0.30, 6.85, 12.70, 0.38, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CLASSIFICATION SETUP
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Model 2: Random Forest Classification", "Predicting whether a crime incident will lead to an arrest")
footer(sl)

add_text(sl, "Problem", 0.35, 1.20, 6.0, 0.42, size=16, bold=True, color=BLUE)
add_text(sl, "Given the characteristics of a reported crime — type, location, time, district — "
             "can we predict whether it will result in an ARREST?",
         0.35, 1.68, 6.0, 0.90, size=13, color=DARK)

add_text(sl, "Why Random Forest?", 0.35, 2.75, 6.0, 0.42, size=16, bold=True, color=BLUE)
add_bullets(sl, [
    "Handles mixed features natively (numeric, binary, categorical)",
    "class_weight='balanced' compensates for 6:1 class imbalance",
    "Feature importance scores reveal what drives arrest outcomes",
    "More robust than a single Decision Tree — reduces overfitting",
], 0.35, 3.22, 6.0, 2.40, size=13, spacing_after=70000)

# Config table
add_text(sl, "Configuration", 6.80, 1.20, 6.20, 0.42, size=16, bold=True, color=BLUE)
config = [
    ("n_estimators",     "100 trees"),
    ("max_depth",        "12"),
    ("class_weight",     "'balanced'"),
    ("Train / Test",     "80% / 20%  (stratified)"),
    ("Target variable",  "Arrest  (0 = no arrest, 1 = arrested)"),
    ("Features used",    "12  (hour, month, day, district, ward, beat, domestic, outdoor, time bin, community area, crime type, city sector)"),
]
for ri, (param, val) in enumerate(config):
    bg = BLUE_LIGHT if ri % 2 == 0 else WHITE
    add_rect(sl, 6.80, 1.72 + ri * 0.70, 2.80, 0.65, bg)
    add_rect(sl, 9.60, 1.72 + ri * 0.70, 3.55, 0.65, bg)
    add_text(sl, param, 6.92, 1.75 + ri * 0.70, 2.65, 0.58,
             size=12, bold=True, color=BLUE)
    add_text(sl, val,   9.72, 1.75 + ri * 0.70, 3.40, 0.58,
             size=11, color=DARK)

add_text(sl, "Class imbalance: 85.8% No Arrest  vs  14.2% Arrest  — balanced weighting is essential",
         0.35, 6.45, 12.60, 0.55, size=13, bold=False, color=GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CONFUSION MATRIX + RESULTS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Model 2: Classification Results", "Random Forest — Accuracy 84.3%  |  ROC-AUC 0.8154")
footer(sl)

add_image(sl, "confusion_matrix.png", 0.25, 1.12, 6.0)

# Results table
results = [
    ("Overall Accuracy",    "84.29%",  "43,434 of 51,510 correct"),
    ("ROC-AUC Score",       "0.8154",  "Strong discrimination power"),
    ("No Arrest Precision", "0.92",    "92% of predicted no-arrests are correct"),
    ("Arrest Recall",       "0.55",    "55% of actual arrests caught"),
    ("Arrest F1-Score",     "0.49",    "Balanced precision-recall tradeoff"),
]
add_text(sl, "Performance Metrics", 6.55, 1.20, 6.50, 0.45, size=16, bold=True, color=BLUE)
for ri, (metric, val, note) in enumerate(results):
    bg = BLUE_LIGHT if ri % 2 == 0 else WHITE
    add_rect(sl, 6.55, 1.75 + ri * 0.72, 6.50, 0.66, bg)
    add_text(sl, metric, 6.65, 1.77 + ri * 0.72, 2.60, 0.30, size=12, bold=True, color=BLUE)
    add_text(sl, val,    9.25, 1.77 + ri * 0.72, 1.00, 0.30, size=14, bold=True, color=ACCENT)
    add_text(sl, note,   6.65, 2.08 + ri * 0.72, 6.30, 0.28, size=10, color=GRAY, italic=True)

add_rect(sl, 6.55, 5.45, 6.50, 0.75, RGBColor(0xFF, 0xF3, 0xE0))
add_text(sl, "An ROC-AUC of 0.8154 means the model is dramatically better than a random baseline (0.5) "
             "at separating arrest from non-arrest cases.",
         6.65, 5.49, 6.30, 0.65, size=12, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — FEATURE IMPORTANCE + ROC
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Model 2: What Drives Arrest Outcomes?", "Feature Importance & ROC Curve")
footer(sl)

add_image(sl, "feature_importance.png", 0.25, 1.12, 7.60)
add_image(sl, "roc_curve.png",          8.10, 1.12, 5.00)

add_text(sl, "Crime type is the #1 predictor — what you do determines whether you get arrested, more than where or when.",
         0.30, 6.82, 12.70, 0.40, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — KEY FINDINGS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Key Findings & Insights", "Data-driven discoveries across all three analyses")
footer(sl)

sections = [
    ("Geographic", BLUE, [
        "Fuller Park: 243 crimes/1k residents — 10× more than Edison Park (24/1k)",
        "10 K-Means clusters reveal distinct hotspot zones across the city",
        "Central/West sector has the highest raw crime count (94,060 incidents, 36.5%)",
    ]),
    ("Temporal", RGBColor(0x2E, 0x7D, 0x32), [
        "Afternoon (31.2%) + Evening (28.5%) = nearly 60% of all incidents",
        "Evening outdoor clusters have the highest arrest rates (~20%)",
        "Late-night indoor crimes (2 AM avg) have the lowest arrest rates (10-11%)",
    ]),
    ("Predictive", ACCENT, [
        "Crime type is the #1 predictor of arrest — stronger than location or time",
        "Community Area & Beat are 2nd and 3rd — local enforcement capacity matters",
        "84.3% accuracy and AUC=0.8154 with no external data beyond the crime record",
    ]),
]
for col_i, (title, color, points) in enumerate(sections):
    bx = 0.30 + col_i * 4.36
    add_rect(sl, bx, 1.18, 4.22, 5.80, RGBColor(0xF5, 0xF8, 0xFF))
    add_rect(sl, bx, 1.18, 4.22, 0.45, color)
    add_text(sl, title, bx + 0.15, 1.21, 4.0, 0.40, size=15, bold=True, color=WHITE)
    add_bullets(sl, points, bx + 0.18, 1.72, 3.90, 5.10,
                size=13, color=DARK, indent_color=color, spacing_after=80000)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — FUTURE WORK
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Limitations & Future Work", "Honest assessment and actionable next steps")
footer(sl)

add_text(sl, "Current Limitations", 0.35, 1.20, 6.20, 0.45, size=16, bold=True, color=ACCENT)
add_bullets(sl, [
    ("Class imbalance", "Only 13.85% arrests → recall = 0.55 even with balanced weights. SMOTE or cost-sensitive learning could improve minority-class detection."),
    ("Population proxy", "Residential Census data underestimates exposure in commercial areas like the Loop (large daytime population not captured)."),
    ("Single year", "2024 data only — multi-year analysis needed to detect trends and pre/post-COVID shifts."),
    ("K-Means geometry", "Assumes spherical clusters — may not fit Chicago's elongated north-south geography well."),
], 0.35, 1.72, 6.20, 5.10, size=12, indent_color=ACCENT, spacing_after=90000)

add_text(sl, "Future Improvements", 6.85, 1.20, 6.15, 0.45, size=16, bold=True, color=BLUE)
add_bullets(sl, [
    ("Frequent Pattern Mining", "Apply Apriori to find association rules: e.g., 'Theft + Street + Afternoon' co-occur 3× more than expected."),
    ("DBSCAN clustering", "Discovers arbitrary-shape clusters and noise points — better suited to Chicago's irregular geography."),
    ("Weather integration", "Merge NOAA daily temperature data to test the warm-weather crime hypothesis."),
    ("Live prediction dashboard", "Extend the D3.js dashboard with real-time arrest probability scores from the trained model."),
], 6.85, 1.72, 6.15, 5.10, size=12, indent_color=BLUE, spacing_after=90000)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CONCLUSION + GITHUB
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.50, BLUE)
add_rect(sl, 0, 4.60, 13.33, 2.90, RGBColor(0x14, 0x36, 0x78))
add_rect(sl, 0, 4.57, 13.33, 0.06, ACCENT)

add_text(sl, "Conclusion", 0.6, 0.45, 12.0, 0.70,
         size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

deliverables = [
    ("99.4%",  "records retained\nafter cleaning"),
    ("84.21%", "variance captured\nby 2 PCA components"),
    ("77",     "community areas\nintegrated"),
    ("K=10",   "clusters validated\nby Elbow Method"),
    ("84.3%",  "classification\naccuracy"),
    ("0.8154", "ROC-AUC score\n(arrest prediction)"),
]
for i, (val, lbl) in enumerate(deliverables):
    bx = 0.40 + (i % 3) * 4.18
    by = 1.30 + (i // 3) * 1.60
    add_rect(sl, bx, by, 3.85, 1.40, RGBColor(0x1A, 0x3E, 0x8A))
    add_text(sl, val, bx + 0.10, by + 0.08, 3.65, 0.65,
             size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(sl, lbl, bx + 0.10, by + 0.72, 3.65, 0.60,
             size=12, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

add_text(sl, "github.com/samagannurdinov",
         0.6, 4.72, 12.0, 0.50, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Data Mining  ·  April 2026  ·  Sama Gannurdinov",
         0.6, 5.28, 12.0, 0.45, size=14, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
add_text(sl, "Thank You",
         0.6, 5.90, 12.0, 0.80, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(sl, 4.5, 6.78, 4.33, 0.05, ACCENT)


# ── Save ───────────────────────────────────────────────────────────────────────
out = "Presentation.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
